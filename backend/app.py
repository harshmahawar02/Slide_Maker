from flask import Flask, request, send_file, jsonify
from flask_cors import CORS
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor
import io
import traceback
import os

app = Flask(__name__)
CORS(app, expose_headers='Content-Disposition')

# Configuration
MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB
MAX_IMAGE_SIZE = 10 * 1024 * 1024  # 10MB
ALLOWED_IMAGE_EXTENSIONS = {'png', 'jpg', 'jpeg', 'gif', 'bmp'}

# Predefined layouts that will be matched against the uploaded PPTX
PREDEFINED_LAYOUTS = {
    'title_content': {
        'name': 'Title and Content',
        'keywords': ['title', 'content', 'body', 'text', 'only'],
        'has_image': False,
        'content_boxes': 1
    },
    'title_two_content': {
        'name': 'Title and Two Content',
        'keywords': ['two', 'comparison', 'columns', 'divided'],
        'has_image': False,
        'content_boxes': 2
    },
    'title_image_content': {
        'name': 'Title Image and Content',
        'keywords': ['picture', 'image', 'content', 'photo', 'text'],
        'has_image': True,
        'content_boxes': 1
    },
    'title_image': {
        'name': 'Title and Image',
        'keywords': ['picture', 'image', 'photo', 'blank'],
        'has_image': True,
        'content_boxes': 0,
        'prefer_simple': True  # Prefer simpler layouts
    }
}

def validate_file_size(file_storage, max_size, file_type="File"):
    """Validate file size before processing"""
    file_storage.seek(0, os.SEEK_END)
    size = file_storage.tell()
    file_storage.seek(0)
    
    if size > max_size:
        size_mb = size / (1024 * 1024)
        max_mb = max_size / (1024 * 1024)
        raise ValueError(f"{file_type} too large: {size_mb:.2f}MB (max: {max_mb}MB)")
    
    return size

def validate_image_format(filename):
    """Validate image file extension"""
    if '.' not in filename:
        return False
    ext = filename.rsplit('.', 1)[1].lower()
    return ext in ALLOWED_IMAGE_EXTENSIONS

def find_layout_by_type(slide_master, layout_type):
    """
    Find the best matching layout based on predefined layout type.
    Uses keyword matching and structure analysis.
    Excludes decorative/branded layouts like 'Mango', 'Cover', etc.
    """
    if layout_type not in PREDEFINED_LAYOUTS:
        return None
    
    layout_config = PREDEFINED_LAYOUTS[layout_type]
    keywords = layout_config['keywords']
    
    # Exclude patterns - layouts to avoid
    exclude_patterns = ['mango', 'cover', 'branded', 'anvil', 'thank', 'section', 'divider']
    
    # First try: Exact name match (excluding branded layouts)
    for layout in slide_master.slide_layouts:
        layout_name_lower = layout.name.lower()
        
        # Skip if matches exclude pattern
        if any(pattern in layout_name_lower for pattern in exclude_patterns):
            continue
            
        if layout_config['name'].lower() == layout_name_lower:
            print(f"✓ Exact match found: {layout.name}")
            return layout
    
    # Second try: Keyword match (excluding branded layouts)
    best_match = None
    best_score = 0
    
    for layout in slide_master.slide_layouts:
        layout_name_lower = layout.name.lower()
        
        # Skip if matches exclude pattern
        if any(pattern in layout_name_lower for pattern in exclude_patterns):
            continue
        
        score = 0
        
        # Count keyword matches
        for keyword in keywords:
            if keyword in layout_name_lower:
                score += 1
        
        if score > best_score:
            best_score = score
            best_match = layout
    
    if best_match and best_score > 0:
        print(f"✓ Keyword match found: {best_match.name} (score: {best_score})")
        return best_match
    
    # Third try: Structure-based matching (excluding branded layouts)
    best_structure_match = None
    best_structure_score = 0
    
    for layout in slide_master.slide_layouts:
        layout_name_lower = layout.name.lower()
        
        # Skip if matches exclude pattern
        if any(pattern in layout_name_lower for pattern in exclude_patterns):
            continue
        
        has_title = False
        body_count = 0
        has_picture = False
        placeholder_count = 0
        
        for shape in layout.placeholders:
            placeholder_count += 1
            try:
                phf = shape.placeholder_format
                if phf.type == PP_PLACEHOLDER.TITLE:
                    has_title = True
                elif phf.type in [PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT]:
                    body_count += 1
                elif phf.type == PP_PLACEHOLDER.PICTURE:
                    has_picture = True
            except:
                continue
        
        # Calculate structure score (prefer simpler layouts with fewer placeholders)
        structure_score = 0
        
        # Match based on structure
        if layout_type == 'title_content':
            if has_title and body_count >= 1 and not has_picture:
                structure_score = 100 - placeholder_count  # Prefer fewer placeholders
        elif layout_type == 'title_two_content':
            if has_title and body_count >= 2:
                structure_score = 100 - placeholder_count
        elif layout_type == 'title_image_content':
            if has_title and body_count >= 1 and has_picture:
                structure_score = 100 - placeholder_count
        elif layout_type == 'title_image':
            if has_title and has_picture and body_count == 0:
                # For title+image, strongly prefer layouts with NO body placeholders
                structure_score = 200 - placeholder_count
            elif has_title and has_picture:
                # Accept layouts with body, but score them lower
                structure_score = 50 - placeholder_count
        
        if structure_score > best_structure_score:
            best_structure_score = structure_score
            best_structure_match = layout
    
    if best_structure_match:
        print(f"✓ Structure match found: {best_structure_match.name} (score: {best_structure_score})")
        return best_structure_match
    
    # Fallback: Use the most appropriate default layout
    print(f"! No perfect match, using fallback for {layout_type}")
    for layout in slide_master.slide_layouts:
        has_title = False
        has_body = False
        
        for shape in layout.placeholders:
            try:
                phf = shape.placeholder_format
                if phf.type == PP_PLACEHOLDER.TITLE:
                    has_title = True
                elif phf.type in [PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT]:
                    has_body = True
            except:
                continue
        
        if has_title and has_body:
            print(f"✓ Fallback layout: {layout.name}")
            return layout
    
    # Ultimate fallback: return first non-blank layout
    if len(slide_master.slide_layouts) > 1:
        print(f"✓ Using second layout as ultimate fallback")
        return slide_master.slide_layouts[1]
    
    return slide_master.slide_layouts[0] if len(slide_master.slide_layouts) > 0 else None

@app.route('/api/debug-layouts', methods=['POST'])
def debug_layouts():
    """Debug endpoint to see all available layouts in the uploaded template"""
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No PowerPoint file uploaded'}), 400
        
        file = request.files['file']
        
        if file.filename == '':
            return jsonify({'error': 'No file selected'}), 400
        
        try:
            prs = Presentation(file)
        except Exception as e:
            return jsonify({'error': f'Invalid PowerPoint file: {str(e)}'}), 400
        
        slide_master = prs.slide_masters[0]
        layouts_info = []
        
        for idx, layout in enumerate(slide_master.slide_layouts):
            placeholders = []
            for shape in layout.placeholders:
                try:
                    phf = shape.placeholder_format
                    placeholders.append({
                        'type': str(phf.type),
                        'name': shape.name
                    })
                except:
                    pass
            
            layouts_info.append({
                'index': idx,
                'name': layout.name,
                'placeholders': placeholders
            })
        
        return jsonify({
            'total_layouts': len(slide_master.slide_layouts),
            'layouts': layouts_info
        }), 200
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/add-slide', methods=['POST'])
def add_slide():
    try:
        # Validate file upload
        if 'file' not in request.files:
            return jsonify({'error': 'No PowerPoint file uploaded'}), 400
        
        file = request.files['file']
        
        if file.filename == '':
            return jsonify({'error': 'No file selected'}), 400
        
        if not file.filename.endswith('.pptx'):
            return jsonify({'error': 'File must be a .pptx PowerPoint file'}), 400
        
        # Validate file size
        try:
            validate_file_size(file, MAX_FILE_SIZE, "PowerPoint file")
        except ValueError as e:
            return jsonify({'error': str(e)}), 400
            
        # Get form data
        layout_type = request.form.get('layout', '').strip()
        title = request.form.get('title', '').strip()
        text = request.form.get('text', '').strip()
        text2 = request.form.get('text2', '').strip()  # For two-content layout
        image = request.files.get('image')
        
        # Validate layout type provided
        if not layout_type or layout_type not in PREDEFINED_LAYOUTS:
            return jsonify({'error': 'Valid layout type is required'}), 400
        
        layout_config = PREDEFINED_LAYOUTS[layout_type]
        
        # Validate required content based on layout
        if layout_config['content_boxes'] > 0 and not text:
            return jsonify({'error': 'Content text is required for this layout'}), 400
        
        if layout_config['content_boxes'] == 2 and not text2:
            return jsonify({'error': 'Second content box is required for this layout'}), 400
        
        # Validate image if provided
        if image:
            if not validate_image_format(image.filename):
                return jsonify({'error': f'Invalid image format. Allowed: {", ".join(ALLOWED_IMAGE_EXTENSIONS)}'}), 400
            try:
                validate_file_size(image, MAX_IMAGE_SIZE, "Image")
            except ValueError as e:
                return jsonify({'error': str(e)}), 400
        
        # Validate position
        try:
            position = int(request.form.get('position', 0))
            if position < 0:
                return jsonify({'error': 'Position cannot be negative'}), 400
        except ValueError:
            return jsonify({'error': 'Position must be a valid number'}), 400

        # Load presentation
        try:
            prs = Presentation(file)
        except Exception as e:
            return jsonify({'error': f'Invalid or corrupted PowerPoint file: {str(e)}'}), 400

        # Find the layout by type
        slide_master = prs.slide_masters[0]
        
        if len(slide_master.slide_layouts) == 0:
            return jsonify({'error': 'Presentation has no slide layouts'}), 400
        
        slide_layout = find_layout_by_type(slide_master, layout_type)
        
        if not slide_layout:
            return jsonify({'error': f'Could not find suitable layout for: {layout_config["name"]}'}), 400

        print(f"\n=== Adding Slide ===")
        print(f"Requested Layout Type: {layout_type}")
        print(f"Using Layout: {slide_layout.name}")
        print(f"Title: {title}")
        print(f"Text: {text[:50]}..." if text and len(text) > 50 else f"Text: {text}")
        if text2:
            print(f"Text2: {text2[:50]}..." if len(text2) > 50 else f"Text2: {text2}")

        # Validate and adjust position
        total_slides = len(prs.slides)
        if position < 0:
            position = 0
        elif position > total_slides:
            position = total_slides

        # Create new slide
        new_slide = prs.slides.add_slide(slide_layout)
        
        # Remove background to make it plain white
        try:
            background = new_slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background
            print(f"✓ Background set to plain white")
        except Exception as e:
            print(f"! Could not modify background: {str(e)}")
        
        # Move slide to correct position if needed
        if position < total_slides:
            try:
                # Use private API carefully - wrap in try/except for safety
                xml_slides = prs.slides._sldIdLst
                slides = list(xml_slides)
                
                # Validate we have the slide to move
                if len(slides) > 0:
                    moved_slide = slides.pop()
                    # Ensure position is within valid range
                    insert_pos = min(position, len(slides))
                    slides.insert(insert_pos, moved_slide)
                    xml_slides.clear()
                    for slide in slides:
                        xml_slides.append(slide)
                    print(f"Slide inserted at position {insert_pos}")
                else:
                    print(f"! No slides to reposition")
            except AttributeError:
                print(f"! Slide positioning not supported in this version of python-pptx")
            except Exception as e:
                print(f"! Could not reposition slide: {str(e)}")
                # Don't fail the entire operation if positioning fails

        # SET TITLE
        title_set = False
        if title:
            try:
                if new_slide.shapes.title:
                    new_slide.shapes.title.text = title
                    print(f"✓ Title: '{title}'")
                    title_set = True
            except Exception as e:
                print(f"! Could not set title: {str(e)}")

        # SET CONTENT - Handle different layout types
        if layout_type == 'title_two_content' and text and text2:
            # Handle two content boxes
            content_placeholders = []
            for shape in new_slide.placeholders:
                try:
                    phf = shape.placeholder_format
                    if phf.type in [PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT]:
                        content_placeholders.append(shape)
                except:
                    continue
            
            if len(content_placeholders) >= 2:
                # Use first two placeholders
                try:
                    tf1 = content_placeholders[0].text_frame
                    tf1.clear()
                    p1 = tf1.paragraphs[0]
                    p1.text = text
                    p1.alignment = PP_ALIGN.LEFT
                    print(f"✓ Content 1: '{text[:50]}...'")
                    
                    tf2 = content_placeholders[1].text_frame
                    tf2.clear()
                    p2 = tf2.paragraphs[0]
                    p2.text = text2
                    p2.alignment = PP_ALIGN.LEFT
                    print(f"✓ Content 2: '{text2[:50]}...'")
                except Exception as e:
                    print(f"! Error setting two-content: {str(e)}")
            else:
                # Fallback: Create two textboxes
                try:
                    # Left textbox
                    txBox1 = new_slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4.2), Inches(5))
                    tf1 = txBox1.text_frame
                    p1 = tf1.paragraphs[0]
                    p1.text = text
                    p1.alignment = PP_ALIGN.LEFT
                    
                    # Right textbox
                    txBox2 = new_slide.shapes.add_textbox(Inches(5.2), Inches(1.8), Inches(4.2), Inches(5))
                    tf2 = txBox2.text_frame
                    p2 = tf2.paragraphs[0]
                    p2.text = text2
                    p2.alignment = PP_ALIGN.LEFT
                    print(f"✓ Two content boxes created (fallback)")
                except Exception as e:
                    print(f"! Error creating two textboxes: {str(e)}")
        
        elif text:
            # Handle single content box
            content_added = False
            placeholder_found = False
            
            for shape in new_slide.placeholders:
                try:
                    phf = shape.placeholder_format
                    placeholder_found = True
                    
                    # Look for body/content placeholder
                    if phf.type in [PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT]:
                        # Found content placeholder - use it
                        tf = shape.text_frame
                        
                        # Validate text frame is writable
                        if not tf:
                            print(f"! Placeholder has no text frame")
                            continue
                        
                        tf.clear()
                        
                        # Add text with proper formatting
                        p = tf.paragraphs[0]
                        p.text = text
                        p.alignment = PP_ALIGN.LEFT
                        p.level = 0
                        
                        # Adjust text frame margins for better positioning
                        tf.margin_left = Inches(0.1)
                        tf.margin_top = Inches(0.1)
                        tf.word_wrap = True
                        
                        print(f"✓ Content: '{text[:50]}...' (in placeholder)")
                        content_added = True
                        break
                except AttributeError as e:
                    print(f"! Placeholder attribute error: {str(e)}")
                    continue
                except Exception as e:
                    print(f"! Error setting placeholder content: {str(e)}")
                    continue
            
            if not placeholder_found:
                print(f"No placeholders found in slide layout")
            
            # If no placeholder found or failed, create text box
            if not content_added:
                try:
                    left = Inches(0.5)
                    top = Inches(1.8)
                    width = Inches(9)
                    height = Inches(5)
                    
                    txBox = new_slide.shapes.add_textbox(left, top, width, height)
                    tf = txBox.text_frame
                    p = tf.paragraphs[0]
                    p.text = text
                    p.alignment = PP_ALIGN.LEFT
                    
                    tf.word_wrap = True
                    tf.margin_left = Inches(0)
                    tf.margin_top = Inches(0)
                    
                    print(f"✓ Content: '{text[:50]}...' (in textbox fallback)")
                except Exception as e:
                    print(f"! Error creating textbox: {str(e)}")
                    raise ValueError("Could not add content to slide")

        # Add image if provided
        if image:
            image_added = False
            
            # Try to find image placeholder first
            for shape in new_slide.placeholders:
                try:
                    phf = shape.placeholder_format
                    if phf.type == PP_PLACEHOLDER.PICTURE:
                        # Found picture placeholder - use correct method
                        image.seek(0)
                        image_stream = io.BytesIO(image.read())
                        
                        # For picture placeholders, we need to delete it and add image in its place
                        left = shape.left
                        top = shape.top
                        width = shape.width
                        height = shape.height
                        
                        # Remove the placeholder
                        sp = shape.element
                        sp.getparent().remove(sp)
                        
                        # Add the picture in the same position
                        new_slide.shapes.add_picture(image_stream, left, top, width, height)
                        
                        print(f"✓ Image added to picture placeholder position")
                        image_added = True
                        break
                except Exception as e:
                    print(f"! Could not use picture placeholder: {str(e)}")
                    continue
            
            # If no placeholder, add as regular image on the right side
            if not image_added:
                try:
                    image.seek(0)  # Reset file pointer
                    image_stream = io.BytesIO(image.read())
                    
                    # Position image on the right side of the slide
                    left = Inches(6.5)
                    top = Inches(1.5)
                    max_height = Inches(5)
                    max_width = Inches(3)
                    
                    # Add image with size constraints
                    pic = new_slide.shapes.add_picture(image_stream, left, top, height=max_height)
                    
                    # Ensure width doesn't exceed bounds
                    if pic.width > max_width:
                        aspect_ratio = pic.height / pic.width
                        pic.width = max_width
                        pic.height = int(max_width * aspect_ratio)
                    
                    print(f"✓ Image added as shape on right side")
                except Exception as e:
                    print(f"! Image error: {str(e)}")
                    # Don't fail the entire operation if image fails
                    pass

        print(f"=== Complete ===\n")
        
        # Save and return
        output = io.BytesIO()
        try:
            prs.save(output)
            output.seek(0)
        except Exception as e:
            print(f"ERROR saving presentation: {str(e)}")
            return jsonify({'error': f'Failed to save presentation: {str(e)}'}), 500

        # Create a new filename for the modified presentation
        base_name = file.filename.rsplit('.', 1)[0]
        new_filename = f"{base_name}_updated.pptx"

        return send_file(
            output, 
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True,
            download_name=new_filename
        )
    
    except ValueError as e:
        # Validation errors
        return jsonify({'error': str(e)}), 400
        
    except Exception as e:
        print(f"ERROR: {str(e)}")
        print(traceback.format_exc())
        return jsonify({'error': f'Server error: {str(e)}'}), 500

@app.route('/api/health', methods=['GET'])
def health_check():
    return jsonify({'status': 'healthy', 'message': 'SlideMaker API is running'}), 200

@app.route('/api/get-slide-count', methods=['POST'])
def get_slide_count():
    """Get the number of slides in uploaded presentation"""
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No PowerPoint file uploaded'}), 400
        
        file = request.files['file']
        
        if file.filename == '':
            return jsonify({'error': 'No file selected'}), 400
        
        if not file.filename.endswith('.pptx'):
            return jsonify({'error': 'File must be a .pptx PowerPoint file'}), 400
        
        # Load presentation
        try:
            prs = Presentation(file)
        except Exception as e:
            return jsonify({'error': f'Invalid PowerPoint file: {str(e)}'}), 400
        
        return jsonify({
            'total_slides': len(prs.slides),
            'filename': file.filename
        }), 200
        
    except Exception as e:
        print(f"Error in get_slide_count: {str(e)}")
        print(traceback.format_exc())
        return jsonify({'error': f'Server error: {str(e)}'}), 500

if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=5000)
