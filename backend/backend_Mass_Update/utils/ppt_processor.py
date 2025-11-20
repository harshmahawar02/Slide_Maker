import re
from pptx import Presentation


def update_ppt_text(ppt_path, replacements):
    """
    Updates all occurrences of specified words in the PPT file (including tables).
    Handles replacements across run boundaries while preserving run count and styles
    (new text is assigned to the first run; subsequent runs are cleared to keep style info).
    Returns True if any replacement was made, False otherwise.
    """
    prs = Presentation(ppt_path)
    updated = False
    # Precompile regex patterns for all replacements
    patterns = [(re.compile(re.escape(old)), new) for old, new in replacements.items()]
    keys = set(replacements.keys())

    def process_paragraph(paragraph):
        """Replace across the concatenated text of runs in a paragraph/cell."""
        nonlocal updated
        if not paragraph.runs:
            return
        original_text = "".join(run.text for run in paragraph.runs)
        # Fast substring check before regex to skip paragraphs without any keys
        if not any(k in original_text for k in keys):
            return
        new_text = original_text
        for pattern, new in patterns:
            new_text = pattern.sub(new, new_text)
        if new_text != original_text:
            # Assign new text to first run and clear others to preserve styles as much as possible
            paragraph.runs[0].text = new_text
            for run in paragraph.runs[1:]:
                run.text = ""
            updated = True

    for slide in prs.slides:
        for shape in slide.shapes:
            # Normal text frames
            if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    process_paragraph(paragraph)
            # Table cells
            table = getattr(shape, "table", None)
            if table is not None:
                for row in table.rows:
                    for cell in row.cells:
                        tf = getattr(cell, "text_frame", None)
                        if tf is not None:
                            for paragraph in tf.paragraphs:
                                process_paragraph(paragraph)

    if updated:
        prs.save(ppt_path)
    return updated
