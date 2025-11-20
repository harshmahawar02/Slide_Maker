import os
import tempfile
import unittest
from pptx import Presentation
from pptx.util import Pt
from utils.ppt_processor import update_ppt_text

def make_test_pptx_with_table(path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    # Add a text box with split runs
    txBox = slide.shapes.add_textbox(100, 100, 500, 100)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    run1 = p.add_run()
    run1.text = 'Hello '
    run1.font.bold = True
    run2 = p.add_run()
    run2.text = 'World'
    run2.font.italic = True
    # Add a table
    table = slide.shapes.add_table(2, 2, 100, 250, 400, 100).table
    cell = table.cell(0, 0)
    p2 = cell.text_frame.paragraphs[0]
    run3 = p2.add_run()
    run3.text = 'Foo'
    run3.font.size = Pt(20)
    run4 = p2.add_run()
    run4.text = 'Bar'
    run4.font.bold = True
    prs.save(path)

class TestPptProcessor(unittest.TestCase):
    def test_run_preserving_replace(self):
        with tempfile.TemporaryDirectory() as tmp:
            pptx_path = os.path.join(tmp, 'test.pptx')
            make_test_pptx_with_table(pptx_path)
            # Replace 'Hello World' (across runs) with 'Hi All' (should preserve bold/italic pattern)
            replacements = {'Hello World': 'Hi All', 'FooBar': 'BazQux'}
            updated = update_ppt_text(pptx_path, replacements)
            self.assertTrue(updated)
            prs = Presentation(pptx_path)
            # Check text box
            found = False
            for shape in prs.slides[0].shapes:
                if not shape.has_text_frame:
                    continue
                for p in shape.text_frame.paragraphs:
                    text = ''.join(run.text for run in p.runs)
                    if text == 'Hi All':
                        found = True
                        self.assertEqual(len(p.runs), 2)  # Should be split into two runs
                        self.assertEqual(p.runs[0].font.bold, True)
                        self.assertEqual(p.runs[1].font.italic, True)
            self.assertTrue(found)
            # Check table cell
            found = False
            for shape in prs.slides[0].shapes:
                if hasattr(shape, 'table') and shape.table is not None:
                    cell = shape.table.cell(0, 0)
                    p = cell.text_frame.paragraphs[0]
                    text = ''.join(run.text for run in p.runs)
                    if text == 'BazQux':
                        found = True
                        self.assertEqual(len(p.runs), 2)
                        self.assertEqual(p.runs[0].font.size, Pt(20))
                        self.assertEqual(p.runs[1].font.bold, True)
            self.assertTrue(found)

if __name__ == '__main__':
    unittest.main()
